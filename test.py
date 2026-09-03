from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # 제목 슬라이드 레이아웃
    title_slide_layout = prs.slide_layouts[0]
    # 내용 슬라이드 레이아웃 (제목 + 내용)
    bullet_slide_layout = prs.slide_layouts[1]

    # --- 슬라이드 1: 표지 ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "미래 공학반 프로젝트 계획안"
    subtitle.text = "발표자: 전인고 이상우\n날짜: 2025.12.03"

    # --- 슬라이드 2: 목차 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "목차"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "1. 소지품 검사 및 유해 물질 감지"
    p = tf.add_paragraph()
    p.text = "2. 웹캠 이용 이상(불건전) 행동 탐지"
    p = tf.add_paragraph()
    p.text = "3. 택배 발송 완료 알림 시스템"

    # --- 슬라이드 3: 소지품 검사 기기 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "1. 소지품 검사 및 유해 물질 감지"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "개요"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "교내 반입 금지 물품(술, 담배 등) 사전 차단 및 안전 확보"
    
    p = tf.add_paragraph()
    p.text = "구현 방법 (현실적 접근)"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "X-Ray: 개인 구매 불가 및 규제 문제로 현실적 어려움 존재"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "대안: 화학 센서(알코올 센서 등)를 활용한 공기질 분석"
    p = tf.add_paragraph()
    p.level = 2
    p.text = "물품 직접 검색은 어렵으나 실내 흡연/음주 행위 즉각 탐지 가능"
    
    p = tf.add_paragraph()
    p.text = "설치 장소"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "교실 및 기숙사 실내"

    # --- 슬라이드 4: 이상 행동 탐지 (개요) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "2. 웹캠 기반 이상 행동 탐지 (개요)"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "개요"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "교내 사각지대 웹캠 설치 -> 딥러닝 기반 불건전 행위 감지"
    
    p = tf.add_paragraph()
    p.text = "구현 프로세스"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "데이터 수집 -> 모델 학습(Training) -> 실시간 추론(Inference)"
    
    p = tf.add_paragraph()
    p.text = "필요 자원"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "학습용: 고성능 GPU 컴퓨팅 자원 필요"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "추론용: 엣지 디바이스(Edge Device) 필요"

    # --- 슬라이드 5: 이상 행동 탐지 (하드웨어/비용) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "2. 웹캠 기반 이상 행동 탐지 (비용 분석)"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "모델 학습 (Training) 환경"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "구매: NVIDIA Orin NX (약 240만원) / AGX (약 313만원) - 고가"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "임대 (추천): Google Colab (Pro $9.99/월)"
    p = tf.add_paragraph()
    p.level = 2
    p.text = "전략: 무료 버전 테스트 -> 필요 시 유료 결제 (가성비 우수)"

    p = tf.add_paragraph()
    p.text = "실시간 추론 (Inference) 장비"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "추천: NVIDIA Jetson Orin Nano Super (약 50만원)"
    p = tf.add_paragraph()
    p.level = 2
    p.text = "선정 이유: 라즈베리파이 대비 AI 라이브러리 호환성 및 성능 우수"

    # --- 슬라이드 6: 택배 알림 시스템 (개요) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "3. 택배 발송 완료 알림 시스템"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "문제점"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "택배실 관리 부재, 미수령 택배 적재로 인한 공간 부족 및 분실 위험"
    
    p = tf.add_paragraph()
    p.text = "해결 방안"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "OCR/바코드 인식 -> 수령인 자동 식별 -> 알림 발송 시스템"
    
    p = tf.add_paragraph()
    p.text = "데이터베이스(DB) 구축"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "학생 DB: 이름, 전화번호, 알림 신청 여부"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "택배 DB: 운송장 정보 (OCR 인식 데이터)"

    # --- 슬라이드 7: 택배 알림 시스템 (로직) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "3. 시스템 상세 로직 & 단계별 계획"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "개인정보 매칭 알고리즘"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "문제: 운송장 정보 마스킹 (예: 이*우, 010-****-1234)"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "해결: 학생 DB와 교차 검증(Cross-check)하여 대상자 유추"
    
    p = tf.add_paragraph()
    p.text = "단계별 추진 계획"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "1단계: 기초 인프라 (DB 구축 + OCR + 문자/방송 알림)"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "2단계: 고도화 (안면 인식 도입 -> 키오스크 자동 안내)"

    # 파일 저장
    file_name = "미래공학반_프로젝트_계획안.pptx"
    prs.save(file_name)
    print(f"'{file_name}' 파일이 생성되었습니다. 왼쪽 폴더 아이콘을 눌러 다운로드하세요.")

# 함수 실행
create_presentation()